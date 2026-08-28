#!/usr/bin/env python3
"""build_pptx.py -- render a deck spec into a PowerPoint presentation.

This is the *renderer* half of the present-riverware-model skill. It makes no
editorial decisions: a deck-spec JSON says which slides exist, in what order,
and what each one shows; this script draws them from the model's own digest.
What belongs in the deck is the spec author's job, and keeping that judgment
out of the script is what makes the deck reviewable before it is built.

Usage:
    python build_pptx.py model.mdl --spec deck.json          # -> model_deck.pptx
    python build_pptx.py model.mdl --auto                    # generic baseline deck
    python build_pptx.py model.mdl --spec deck.json -o out.pptx
    python build_pptx.py model.mdl --spec deck.json --template client.pptx
    python build_pptx.py model.mdl --spec deck.json --rls policy.rls
    python build_pptx.py model.mdl --spec deck.json --date 2026-08-28

Requires python-pptx:  pip install python-pptx

Spec shape (reference.md holds the authoritative schema):

    {"deck_title": "Arbor Basin operations",
     "date": "2026-08-28",
     "slides": [
       {"type": "title",   "id": "cover",  "notes": "..."},
       {"type": "network", "id": "system", "bullets": ["..."]},
       {"type": "policy",  "id": "rules",
        "refs": {"set": "Arbor Basin Rules", "groups": ["Aspen Rules"]},
        "annotations": {"Aspen Rules": "Guide-curve operation."}},
       {"type": "series",  "id": "pool",
        "refs": {"series": ["Aspen.Pool Elevation", "Birch.Pool Elevation"]}},
       {"type": "caveats", "id": "caveats"}
     ]}

Design notes worth knowing before editing this file:

- The model is never read here. Extraction goes through build_digest() in the
  visualize skill, which wraps the .mdl parser (AGENTS.md hard rule).
- Validation runs to completion before a single slide is drawn. An unknown
  slide type or an unresolvable reference lists the valid targets and exits
  non-zero, so a wrong spec never yields a half-built deck.
- The output is written through a fixed-timestamp zip pass. python-pptx stamps
  each zip entry with the current clock, which would make two runs of the same
  spec differ in bytes while being identical as documents.
- Only PowerPoint validates a .pptx. A clean run here is not proof the deck
  opens; the compatibility matrix in reference.md records what was checked.
"""
from __future__ import annotations

import io
import json
import math
import os
import re
import sys
import zipfile
from datetime import date as _date, datetime, timedelta
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent
                       / "visualize-riverware-model"))
from digest_to_json import (build_digest, layout_nodes,  # noqa: E402
                            policy_from_rls)

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData, XyChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.dml import MSO_LINE_DASH_STYLE, MSO_THEME_COLOR
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, PP_PLACEHOLDER
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Emu, Inches, Pt
    PPTX_ERROR: Exception | None = None
except ImportError as exc:  # reported by main(); validation runs without pptx
    PPTX_ERROR = exc

# ----------------------------------------------------------------- vocabulary
SLIDE_TYPES = ("title", "network", "summary", "policy", "series", "table",
               "chart-xy", "bullets", "caveats")

# ------------------------------------------------------------------ geometry
# Fractions of the slide, so a 4:3 client template lays out like the default
# 16:9 deck instead of spilling off the edge.
SLIDE_W_IN = 13.333        # default slide width, inches (16:9 widescreen)
SLIDE_H_IN = 7.5           # default slide height, inches
MARGIN_FRAC = 0.042        # side margin, 0.02-0.10 of slide width
CONTENT_TOP_FRAC = 0.20    # top of the body area, 0.15-0.30 of slide height
FOOTER_FRAC = 0.075        # strip reserved for the footer, 0.05-0.12 of height

# ------------------------------------------------------------------- palette
# Object-type fills, matched to the dashboard legend (template.html). Order
# matters: the first substring hit wins, so PowerReservoir precedes Reservoir.
TYPE_COLORS = [
    (("PowerReservoir", "InlinePower"), "F6B26B", "Power reservoir / plant"),
    (("Reservoir",), "6FA8DC", "Storage reservoir"),
    (("Reach",), "76C7C0", "Reach"),
    (("Gage",), "B7B7B7", "Stream gage"),
    (("GroundWater",), "C49A6C", "Groundwater"),
    (("WaterUser", "AggDiversionSite"), "93C47D", "Water user"),
    (("Canal",), "8EC9A0", "Canal"),
    (("Diversion",), "FFD966", "Diversion"),
    (("DataObj",), "B4A7D6", "Data object"),
    (("Confluence",), "9FC5E8", "Confluence"),
]
OTHER_FILL, OTHER_LABEL = "D9D9D9", "Other"
NODE_LINE = "5F6B76"       # node outline
FLOW_LINE = "7B8794"       # water-conveyance connector
DATA_LINE = "B8C4CE"       # data/head connector, drawn dashed
# Chart series colors, drawn in order and reused past the sixth series.
SERIES_COLORS = ["2563A8", "76C7C0", "F6B26B", "93C47D", "B4A7D6", "C49A6C"]

# ------------------------------------------------------------ tuning: capacity
MAX_SERIES_POINTS = 400    # points per chart, 50-2000; higher = slower to open
SUMMARY_MAX_ROWS = 14      # objects listed on one summary slide, 5-25
TABLE_MAX_ROWS = 12        # data rows on one table slide, 5-25
POLICY_MAX_ROWS = 18       # group + rule lines on one policy slide, 6-30
BULLET_MAX = 8             # bullets before the text shrinks, 4-12
CAPTION_MAX_SHARE = 0.33   # body area a figure caption may take, 0.15-0.50
NETWORK_CROWDED = 45       # objects past which the schematic is flagged, 20-120
MIN_NODE_FONT_PT = 6.0     # smallest node label, 5-9 pt; below this it is noise
NODE_FONT_RATIO = 0.55     # node label size as a share of box height, 0.3-0.7
NODE_FONT_WARN = 8.0       # labels under this size, in points, are reported
NODE_STRETCH_MAX = 2.5     # tallest a node may get relative to its width, 1-4
MAX_NODE_W_IN = 2.2        # widest a node box may grow, 1.0-3.5 inches
MAX_NODE_H_IN = 0.6        # tallest a node box may grow, 0.3-1.0 inches
LEGEND_ROWS = 2            # rows the object-type key may wrap onto, 1-3

# ------------------------------------------------------------ tuning: text
DESC_CHARS = 90            # object description shown in tables, 40-200 chars
NOTE_CHARS = 110           # spec one-liner shown on a policy row, 40-200 chars

# Fixed document timestamp. Real dates would make two renders of one spec
# differ in bytes; the visible date comes from --date instead.
FIXED_TIMESTAMP = datetime(2001, 1, 1, 0, 0, 0)
FIXED_ZIP_DATE = (1980, 1, 1, 0, 0, 0)

CAVEAT_LINES = [
    "Every figure here comes from the model file itself; nothing is inferred "
    "beyond it.",
    "Result series are the values stored in the model when it was saved, not "
    "a fresh run.",
    "Rule behaviour is summarised. RiverWare's Rule Log is the authority on "
    "what actually fired.",
    "Check any number against RiverWare before it informs a decision.",
]


# ------------------------------------------------------------------- helpers
def fmt_num(v) -> str:
    if v is None:
        return ""
    if isinstance(v, (int, float)):
        if v == int(v) and abs(v) < 1e15:
            return f"{int(v):,}"
        return f"{v:,.4g}" if abs(v) >= 0.001 else f"{v:.3g}"
    return str(v)


def trim(text: str, limit: int) -> str:
    text = (text or "").strip()
    return text if len(text) <= limit else text[:limit - 3].rstrip() + "..."


def color_for(obj_type: str) -> tuple[str, str]:
    for hints, fill, label in TYPE_COLORS:
        if any(h in obj_type for h in hints):
            return fill, label
    return OTHER_FILL, OTHER_LABEL


def iso_day(stamp: str) -> str:
    """RiverWare 'MM-DD-YYYY HH:MM:SS' -> 'YYYY-MM-DD'; anything else unchanged."""
    m = re.match(r"(\d{2})-(\d{2})-(\d{4})", stamp or "")
    return f"{m.group(3)}-{m.group(1)}-{m.group(2)}" if m else (stamp or "?")


def parse_stamp(stamp: str) -> datetime | None:
    m = re.match(r"(\d{2})-(\d{2})-(\d{4})\s+(\d{1,2}):(\d{2})", stamp or "")
    if not m:
        return None
    mo, dd, yy, hh, mi = (int(g) for g in m.groups())
    # RiverWare writes hour 24 for midnight ending a day; timedelta carries it.
    return datetime(yy, mo, dd) + timedelta(hours=hh, minutes=mi)


def advance(when: datetime, count: int, unit: str) -> datetime:
    if unit.startswith("MONTH"):
        month = when.month - 1 + count
        year = when.year + month // 12
        month = month % 12 + 1
        day = min(when.day, [31, 29 if year % 4 == 0 and (year % 100 or not
                                                          year % 400) else 28,
                             31, 30, 31, 30, 31, 31, 30, 31, 30, 31][month - 1])
        return when.replace(year=year, month=month, day=day)
    if unit.startswith("YEAR"):
        return when.replace(year=when.year + count)
    per = {"MINUTE": timedelta(minutes=1), "HOUR": timedelta(hours=1),
           "DAY": timedelta(days=1), "WEEK": timedelta(weeks=1)}
    return when + per.get(unit, timedelta(days=1)) * count


def series_dates(entry: dict) -> list[datetime] | None:
    """Timestamps for one digest series, stepped from its own start and timestep."""
    start = parse_stamp(entry.get("start", ""))
    if start is None:
        return None
    m = re.match(r"(\d+)\s+(\w+)", entry.get("timestep", "") or "")
    count, unit = (int(m.group(1)), m.group(2).upper()) if m else (1, "DAY")
    out, cur = [], start
    for _ in entry["values"]:
        out.append(cur)
        cur = advance(cur, count, unit)
    return out


def align_series(entries: list[dict]) -> tuple[list[datetime], list[list], int]:
    """Put several series on one date axis, gap-filled and thinned to fit.

    Series in a model start and end on different timesteps, so the axis is the
    union of their dates; a series with no value on a date contributes None,
    which PowerPoint draws as a gap rather than a drop to zero. Thinning keeps
    a decade of daily values from making a chart nobody can open.
    """
    per = []
    for e in entries:
        dates = series_dates(e)
        per.append({} if dates is None else dict(zip(dates, e["values"])))
    axis = sorted(set().union(*[set(p) for p in per])) if per else []
    stride = max(1, math.ceil(len(axis) / MAX_SERIES_POINTS)) if axis else 1
    axis = axis[::stride]
    return axis, [[p.get(d) for d in axis] for p in per], stride


# ---------------------------------------------------------------- validation
def _listing(items, limit: int = 12) -> str:
    items = sorted(items)
    shown = ", ".join(items[:limit])
    return shown + (f", ... ({len(items)} total)" if len(items) > limit else "")


def validate_spec(spec: dict, digest: dict) -> list[str]:
    """Every problem with a spec, as messages naming the valid targets.

    Runs before anything is drawn and does not need python-pptx: a spec can be
    checked on a machine that cannot render it.
    """
    errs: list[str] = []
    if not isinstance(spec, dict):
        return [f"spec must be a JSON object, got {type(spec).__name__}"]
    slides = spec.get("slides")
    if not isinstance(slides, list) or not slides:
        return ["spec needs a non-empty 'slides' list"]

    objects = {o["name"]: o for o in digest["objects"]}
    types = {o["type"] for o in digest["objects"]}
    series = {f"{s['object']}.{s['slot']}": s for s in digest["series"]}
    tables = {f"{t['object']}.{t['slot']}": t for t in digest["tables"]}
    sets = {s["name"]: s for s in digest.get("policy", {}).get("sets", [])}
    seen_ids: set[str] = set()

    for i, sl in enumerate(slides, 1):
        if not isinstance(sl, dict):
            errs.append(f"slide {i}: expected an object, got {type(sl).__name__}")
            continue
        kind = sl.get("type")
        where = f"slide {i} ({kind})"
        if kind not in SLIDE_TYPES:
            errs.append(f"slide {i}: type {kind!r} is not one of "
                        f"{', '.join(SLIDE_TYPES)}")
            continue
        sid = sl.get("id")
        if sid is not None:
            if not isinstance(sid, str) or not sid.strip():
                errs.append(f"{where}: 'id' must be a non-empty string")
            elif sid in seen_ids:
                errs.append(f"{where}: duplicate id {sid!r}")
            else:
                seen_ids.add(sid)
        for field in ("bullets",):
            if field in sl and not isinstance(sl[field], list):
                errs.append(f"{where}: '{field}' must be a list of strings")
        refs = sl.get("refs") or {}
        if not isinstance(refs, dict):
            errs.append(f"{where}: 'refs' must be an object")
            continue

        for name in refs.get("objects", []) or []:
            if name not in objects:
                errs.append(f"{where}: object {name!r} is not in this model; "
                            f"available: {_listing(objects)}")
        for name in refs.get("types", []) or []:
            if name not in types:
                errs.append(f"{where}: object type {name!r} is not in this "
                            f"model; available: {_listing(types)}")

        if kind == "series":
            wanted = refs.get("series") or []
            if not wanted:
                errs.append(f"{where}: refs.series must name at least one "
                            f"'Object.Slot' result series")
            have = _listing(series) or "(none: the model stores no results)"
            for ref in wanted:
                if ref not in series:
                    errs.append(
                        f"{where}: {ref!r} is not a stored result series in "
                        f"this model; available: {have}")
        if kind in ("table", "chart-xy"):
            ref = refs.get("table")
            if not ref:
                errs.append(f"{where}: refs.table must name an "
                            f"'Object.Slot' lookup table")
            elif ref not in tables:
                errs.append(f"{where}: {ref!r} is not an extracted lookup "
                            f"table; available: {_listing(tables) or '(none)'}")
            elif kind == "chart-xy":
                cols = tables[ref]["cols"]
                for axis in ("x",):
                    if refs.get(axis) and refs[axis] not in cols:
                        errs.append(f"{where}: column {refs[axis]!r} is not in "
                                    f"{ref}; columns: {_listing(cols)}")
                for col in refs.get("y", []) or []:
                    if col not in cols:
                        errs.append(f"{where}: column {col!r} is not in {ref}; "
                                    f"columns: {_listing(cols)}")
                if len(cols) < 2 and not refs.get("y"):
                    errs.append(f"{where}: {ref} has one column, so there is "
                                f"nothing to plot against; use a 'table' slide")
        if kind == "policy":
            if not sets:
                errs.append(f"{where}: this model carries no rule set. Supply "
                            f"one with --rls, or drop the policy slide")
                continue
            set_name = refs.get("set")
            if set_name and set_name not in sets:
                errs.append(f"{where}: rule set {set_name!r} is not in this "
                            f"model; available: {_listing(sets)}")
                continue
            chosen = sets[set_name] if set_name else next(iter(sets.values()))
            names = {g["name"] for g in chosen["groups"]}
            for g in refs.get("groups", []) or []:
                if g not in names:
                    errs.append(f"{where}: group {g!r} is not in rule set "
                                f"{chosen['name']!r}; available: {_listing(names)}")
        if kind == "bullets" and not sl.get("bullets"):
            errs.append(f"{where}: a bullets slide needs a 'bullets' list")
    return errs


# --------------------------------------------------------------- auto baseline
# Objects and slots the baseline deck leads with when no spec was written.
AUTO_SERIES_SLOTS = ["Pool Elevation", "Outflow"]
AUTO_SERIES_OBJECTS = 4    # reservoirs per baseline chart, 2-8
AUTO_POLICY_GROUPS = 4     # rule groups per baseline policy slide, 2-8


def auto_spec(digest: dict, deck_title: str, when: str) -> dict:
    """A generic baseline deck: orientation, policy skeleton, key series.

    This is the escape hatch for a smoke test or a user working without an AI
    layer. It is deliberately unsurprising -- it makes no claim about which
    part of the model matters, which is exactly what a written spec is for.
    """
    slides: list[dict] = [
        {"type": "title", "id": "cover"},
        {"type": "network", "id": "network"},
        {"type": "summary", "id": "summary"},
    ]
    for st in digest.get("policy", {}).get("sets", []):
        groups = [g for g in st["groups"] if g["kind"] == "POLICY_GROUP"]
        for start in range(0, len(groups), AUTO_POLICY_GROUPS):
            chunk = groups[start:start + AUTO_POLICY_GROUPS]
            slides.append({"type": "policy",
                           "id": f"policy-{len(slides)}",
                           "title": f"Policy: {st['name']}",
                           "refs": {"set": st["name"],
                                    "groups": [g["name"] for g in chunk]}})
    reservoirs = [o["name"] for o in digest["objects"] if "Reservoir" in o["type"]]
    have = {f"{s['object']}.{s['slot']}" for s in digest["series"]}
    for slot in AUTO_SERIES_SLOTS:
        refs = [f"{n}.{slot}" for n in reservoirs if f"{n}.{slot}" in have]
        if refs:
            slides.append({"type": "series", "id": f"series-{slot.lower()}",
                           "title": f"Reservoir {slot.lower()}",
                           "refs": {"series": refs[:AUTO_SERIES_OBJECTS]}})
    slides.append({"type": "caveats", "id": "caveats"})
    return {"deck_title": deck_title, "date": when, "template": None,
            "slides": slides}


# ---------------------------------------------------------------------- deck
# Layout names to look for in a template, and the index to fall back on in the
# stock python-pptx template when a client template names its layouts freely.
LAYOUT_CHOICES = {
    "title": (("title slide",), 0),
    "title_content": (("title and content",), 1),
    "title_only": (("title only",), 5),
    "blank": (("blank",), 6),
}
BODY_PLACEHOLDERS = (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT,
                     PP_PLACEHOLDER.SUBTITLE) if PPTX_ERROR is None else ()


class Deck:
    """A presentation under construction, plus the theme-inheritance guesswork.

    A client template is honoured on a best-effort contract (REQ-007): layouts
    are found by name, then by index, then by giving up and drawing a text box.
    Every fallback is recorded in `warnings` so the review step sees it.
    """

    def __init__(self, template: Path | None, footer: str):
        self.warnings: list[str] = []
        self.footer = footer
        if template is not None:
            self.prs = Presentation(str(template))
            self._drop_template_slides()
        else:
            self.prs = Presentation()
            self.prs.slide_width = Inches(SLIDE_W_IN)
            self.prs.slide_height = Inches(SLIDE_H_IN)
        self.width = int(self.prs.slide_width)
        self.height = int(self.prs.slide_height)

    def _drop_template_slides(self) -> None:
        """Keep the template's masters and theme, discard its example slides."""
        id_list = self.prs.slides._sldIdLst
        for slide_id in list(id_list):
            self.prs.part.drop_rel(slide_id.rId)
            id_list.remove(slide_id)

    def _layout(self, hint: str):
        names, index = LAYOUT_CHOICES[hint]
        layouts = self.prs.slide_layouts
        for layout in layouts:
            if layout.name.strip().lower() in names:
                return layout
        for layout in layouts:
            if any(n in layout.name.strip().lower() for n in names):
                return layout
        if index < len(layouts):
            self.warnings.append(
                f"template has no '{names[0]}' layout; used layout "
                f"{index} ({layouts[index].name!r})")
            return layouts[index]
        self.warnings.append(
            f"template has no '{names[0]}' layout and too few layouts to fall "
            f"back on; used {layouts[0].name!r}")
        return layouts[0]

    # -- geometry ----------------------------------------------------------
    @property
    def margin(self) -> int:
        return int(self.width * MARGIN_FRAC)

    def content_box(self) -> tuple[int, int, int, int]:
        """left, top, width, height of the area below the slide title."""
        top = int(self.height * CONTENT_TOP_FRAC)
        return (self.margin, top, self.width - 2 * self.margin,
                self.height - top - int(self.height * FOOTER_FRAC))

    # -- slides ------------------------------------------------------------
    def add_slide(self, hint: str, title: str | None = None):
        slide = self.prs.slides.add_slide(self._layout(hint))
        if title is not None:
            self.set_title(slide, title)
        return slide

    def set_title(self, slide, text: str) -> None:
        holder = slide.shapes.title
        if holder is None:
            holder = slide.shapes.add_textbox(
                self.margin, int(self.height * 0.045),
                self.width - 2 * self.margin, int(self.height * 0.12))
            holder.text_frame.paragraphs[0].font.size = Pt(28)
            self.warnings.append("a layout has no title placeholder; the slide "
                                 "title was drawn as a plain text box")
        frame = holder.text_frame
        frame.text = text
        frame.word_wrap = True
        if len(text) > 60:
            for para in frame.paragraphs:
                para.font.size = Pt(24)

    def finish_slide(self, slide, notes: str) -> None:
        """Footer, notes, and removal of the placeholders nothing filled."""
        for shape in list(slide.placeholders):
            fmt = shape.placeholder_format
            if fmt.type in BODY_PLACEHOLDERS and not shape.text_frame.text.strip():
                shape._element.getparent().remove(shape._element)
        box = slide.shapes.add_textbox(
            self.margin, self.height - int(self.height * FOOTER_FRAC * 0.85),
            self.width - 2 * self.margin, int(self.height * 0.05))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = self.footer
        run.font.size = Pt(9)
        run.font.color.theme_color = MSO_THEME_COLOR.TEXT_1
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    # -- primitives --------------------------------------------------------
    def text_block(self, slide, left, top, width, height, lines,
                   size=16, bullet=True, color=None):
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.word_wrap = True
        for i, line in enumerate(lines):
            para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
            run = para.add_run()
            run.text = (f"- {line}" if bullet else str(line))
            run.font.size = Pt(size)
            if color is None:
                run.font.color.theme_color = MSO_THEME_COLOR.TEXT_1
            else:
                run.font.color.rgb = RGBColor.from_string(color)
            para.space_after = Pt(size * 0.45)
        return box

    def table(self, slide, left, top, width, height, header, rows, widths=None,
              size=12):
        shape = slide.shapes.add_table(len(rows) + 1, len(header),
                                       left, top, width, height)
        table = shape.table
        if widths:
            for i, frac in enumerate(widths):
                table.columns[i].width = Emu(int(width * frac))
        for c, label in enumerate(header):
            self._cell(table.cell(0, c), label, size, bold=True)
        for r, row in enumerate(rows, 1):
            for c, value in enumerate(row):
                self._cell(table.cell(r, c), value, size)
        return table

    @staticmethod
    def _cell(cell, value, size, bold=False):
        cell.margin_left = Inches(0.06)
        cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        frame = cell.text_frame
        frame.word_wrap = True
        frame.text = "" if value is None else str(value)
        for para in frame.paragraphs:
            para.font.size = Pt(size)
            para.font.bold = bold


def send_to_back(shape) -> None:
    """Move a shape behind its siblings, past the group's own two headers."""
    element = shape._element
    tree = element.getparent()
    tree.remove(element)
    tree.insert(2, element)


def add_arrow_head(line) -> None:
    """Arrowhead on a connector. python-pptx has no API for it, so set the XML.

    a:tailEnd is the last child of a:ln in schema order, so appending is safe
    as long as dash and color are set first.
    """
    ln = line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"),
                             {"type": "triangle", "w": "sm", "len": "sm"}))


# ------------------------------------------------------------ slide renderers
class Renderer:
    """Draws one deck. Every method reads the digest and the spec, nothing else."""

    def __init__(self, deck: Deck, digest: dict, spec: dict, when: str):
        self.deck = deck
        self.digest = digest
        self.spec = spec
        self.when = when
        self.model_file = digest["model"]["file"]
        self.series = {f"{s['object']}.{s['slot']}": s for s in digest["series"]}
        self.tables = {f"{t['object']}.{t['slot']}": t for t in digest["tables"]}
        self.sets = {s["name"]: s
                     for s in digest.get("policy", {}).get("sets", [])}
        self.notes: list[str] = []

    # -- shared ------------------------------------------------------------
    def render(self, slide_spec: dict) -> None:
        kind = slide_spec["type"]
        method = getattr(self, "_" + kind.replace("-", "_"))
        slide, provenance = method(slide_spec)
        note = "\n".join([p for p in provenance if p]
                         + ([slide_spec["notes"]] if slide_spec.get("notes")
                            else []))
        self.deck.finish_slide(slide, note)

    def _title_for(self, slide_spec: dict, default: str) -> str:
        return slide_spec.get("title") or default

    def _bullets_under(self, slide, spec, left, top, width, height):
        """Spec narrative under a figure, with the figure keeping the space.

        The height a block of text needs depends on how it wraps, not on how
        many bullets there are. Guessing low pushes the last line over the
        footer; guessing high starves the figure, so the text is shrunk until
        it fits the share it is allowed.
        """
        lines = [str(b) for b in (spec.get("bullets") or []) if str(b).strip()]
        if not lines:
            return top, height
        budget = height * CAPTION_MAX_SHARE
        for size in (14, 12, 10):
            columns = max(20, int(width / 12700 / (size * 0.5)))
            rows = sum(max(1, math.ceil(len(line) / columns)) for line in lines)
            used = int(Inches(size * 1.45 / 72) * rows
                       + Inches(0.06) * len(lines))
            if used <= budget:
                break
        else:
            used = int(budget)
            self.deck.warnings.append(
                f"{len(lines)} bullets on a figure slide were cut to 10 pt to "
                f"leave room for the figure; consider shortening them")
        self.deck.text_block(slide, left, top + height - used, width, used,
                             lines, size=size)
        return top, height - used

    # -- title -------------------------------------------------------------
    def _title(self, spec):
        deck = self.deck
        slide = deck.add_slide("title")
        model = self.digest["model"]
        deck.set_title(slide, self._title_for(
            spec, self.spec.get("deck_title") or model["file"]))
        run = model.get("run") or {}
        counts = sorted(model["type_counts"].items(), key=lambda kv: (-kv[1], kv[0]))
        lead = ", ".join(f"{n} {t}" for t, n in counts[:3])
        more = f", +{len(counts) - 3} more types" if len(counts) > 3 else ""
        lines = [f"{model['file']}  (RiverWare {model['version']})"]
        if run.get("start"):
            lines.append(f"Run {iso_day(run['start'])} to {iso_day(run['end'])}"
                         f", timestep {run['timestep']}")
        lines.append(f"{model['object_count']} simulation objects: {lead}{more}")
        lines.append(self.when)
        lines += [str(b) for b in (spec.get("bullets") or [])]

        holder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx != 0 and shape.has_text_frame:
                holder = shape
                break
        if holder is None:
            left, top, width, height = deck.content_box()
            deck.text_block(slide, left, top, width, height, lines,
                            size=16, bullet=False)
        else:
            frame = holder.text_frame
            frame.word_wrap = True
            for i, line in enumerate(lines):
                para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
                run_ = para.add_run()
                run_.text = line
                run_.font.size = Pt(16)
        return slide, [f"Source: {self.model_file}."]

    # -- network -----------------------------------------------------------
    def _network(self, spec):
        deck = self.deck
        slide = deck.add_slide("title_only", self._title_for(
            spec, "Model network"))
        refs = spec.get("refs") or {}
        objects = self.digest["objects"]
        if refs.get("objects"):
            keep = set(refs["objects"])
            objects = [o for o in objects if o["name"] in keep]
        names = {o["name"] for o in objects}
        links = [l for l in self.digest["links"]
                 if l["from"] in names and l["to"] in names]
        left, top, width, height = deck.content_box()
        top, height = self._bullets_under(slide, spec, left, top, width, height)

        legend_h = Inches(0.26) * LEGEND_ROWS
        height -= legend_h
        layout = layout_nodes(objects, links)
        if not layout["nodes"]:
            deck.text_block(slide, left, top, width, height,
                            ["No objects to draw."], bullet=False)
            return slide, ["Source: " + self.model_file + "."]

        # The layered layout is usually far wider than it is tall, so scaling
        # it uniformly wastes the lower half of the slide and leaves labels
        # unreadable. The two axes are scaled independently, each capped by
        # its own fit and by an absolute node size, then by each other, so a
        # sparse network gets big boxes instead of a few enormous ones and a
        # dense network still fits.
        scale_x = min(width / layout["width"],
                      Inches(MAX_NODE_W_IN) / layout["node_width"])
        scale_y = min(height / layout["height"],
                      Inches(MAX_NODE_H_IN) / layout["node_height"])
        scale_y = min(scale_y, scale_x * NODE_STRETCH_MAX)
        scale_x = min(scale_x, scale_y * NODE_STRETCH_MAX)
        off_x = left + (width - layout["width"] * scale_x) / 2
        off_y = top + (height - layout["height"] * scale_y) / 2
        node_w = layout["node_width"] * scale_x
        node_h = layout["node_height"] * scale_y
        font_pt = max(MIN_NODE_FONT_PT,
                      min(12.0, node_h / 12700 * NODE_FONT_RATIO))

        shapes = {}
        seen_types: dict[str, str] = {}
        for node in layout["nodes"]:
            fill, label = color_for(node["type"])
            seen_types.setdefault(label, fill)
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Emu(int(off_x + node["x"] * scale_x)),
                Emu(int(off_y + node["y"] * scale_y)),
                Emu(int(node_w)), Emu(int(node_h)))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor.from_string(fill)
            box.line.color.rgb = RGBColor.from_string(NODE_LINE)
            box.line.width = Pt(0.75)
            box.shadow.inherit = False
            frame = box.text_frame
            frame.word_wrap = False
            frame.margin_left = frame.margin_right = Inches(0.02)
            frame.margin_top = frame.margin_bottom = 0
            para = frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            # A label wider than its box overflows onto its neighbours, so it
            # is cut to what the box can hold at the chosen size.
            run.text = trim(node["name"],
                            max(6, int(node_w / 12700 / (font_pt * 0.55))))
            run.font.size = Pt(font_pt)
            run.font.color.rgb = RGBColor.from_string("1F2933")
            shapes[node["name"]] = box

        for edge in layout["edges"]:
            start, end = shapes[edge["from"]], shapes[edge["to"]]
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, start.left, start.top, end.left, end.top)
            conn.begin_connect(start, 3)   # right edge
            conn.end_connect(end, 1)       # left edge
            flow = edge["kind"] == "flow"
            conn.line.color.rgb = RGBColor.from_string(
                FLOW_LINE if flow else DATA_LINE)
            conn.line.width = Pt(1.25 if flow else 1.0)
            if not flow:
                conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
            else:
                add_arrow_head(conn.line)
            send_to_back(conn)

        self._legend(slide, left, top + height, width, legend_h, seen_types)

        if len(layout["nodes"]) > NETWORK_CROWDED or font_pt < NODE_FONT_WARN:
            deck.warnings.append(
                f"network slide: {len(layout['nodes'])} objects at "
                f"{font_pt:.1f} pt labels -- check it is readable when "
                f"projected, or split it with refs.objects")
        flows = sum(1 for e in layout["edges"] if e["kind"] == "flow")
        return slide, [f"Source: {self.model_file}. {len(layout['nodes'])} "
                       f"objects, {flows} flow links and "
                       f"{len(layout['edges']) - flows} data links. Link "
                       f"direction is inferred from which end is an inflow "
                       f"slot; data and head links are dashed."]

    def _legend(self, slide, left, top, width, height, seen_types):
        """Colour key for the object types actually drawn, wrapped over rows."""
        deck = self.deck
        swatch = Inches(0.15)
        row_h = height / LEGEND_ROWS
        x, row = left, 0
        for index, (label, fill) in enumerate(sorted(seen_types.items())):
            text_w = Inches(0.085) * len(label) + Inches(0.14)
            step = swatch * 1.35 + text_w
            if x + step > left + width:
                row += 1
                x = left
                if row >= LEGEND_ROWS:
                    deck.warnings.append(
                        f"network legend: {len(seen_types) - index} object "
                        f"types did not fit and were left out")
                    return
            y = top + row * row_h
            chip = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y + row_h * 0.2)),
                swatch, swatch)
            chip.fill.solid()
            chip.fill.fore_color.rgb = RGBColor.from_string(fill)
            chip.line.color.rgb = RGBColor.from_string(NODE_LINE)
            chip.line.width = Pt(0.5)
            chip.shadow.inherit = False
            box = slide.shapes.add_textbox(Emu(int(x + swatch * 1.35)),
                                           Emu(int(y + row_h * 0.05)),
                                           Emu(int(text_w)), Emu(int(row_h)))
            box.text_frame.word_wrap = False
            run = box.text_frame.paragraphs[0].add_run()
            run.text = label
            run.font.size = Pt(9)
            run.font.color.theme_color = MSO_THEME_COLOR.TEXT_1
            x += step

    # -- summary -----------------------------------------------------------
    def _summary(self, spec):
        deck = self.deck
        slide = deck.add_slide("title_only",
                               self._title_for(spec, "What the model contains"))
        refs = spec.get("refs") or {}
        objects = self.digest["objects"]
        if refs.get("objects"):
            keep = set(refs["objects"])
            objects = [o for o in objects if o["name"] in keep]
        if refs.get("types"):
            keep = set(refs["types"])
            objects = [o for o in objects if o["type"] in keep]
        objects = sorted(objects, key=lambda o: (o["type"], o["name"]))

        left, top, width, height = deck.content_box()
        top, height = self._bullets_under(slide, spec, left, top, width, height)

        counts = sorted(self.digest["model"]["type_counts"].items(),
                        key=lambda kv: (-kv[1], kv[0]))
        lead = [f"{self.digest['model']['object_count']} objects: "
                + ", ".join(f"{n} {t}" for t, n in counts[:6])
                + (f", +{len(counts) - 6} more types" if len(counts) > 6 else "")]
        no_ruleset = not self.sets and refs.get("ruleset_note") is not False
        if no_ruleset:
            lead.append("Ruleset not included: this deck describes the model's "
                        "structure only, not its operating policy.")
        head_h = Inches(0.32) * len(lead)
        deck.text_block(slide, left, top, width, head_h, lead, size=13,
                        bullet=False)
        top += head_h + Inches(0.1)
        height -= head_h + Inches(0.1)

        shown = objects[:SUMMARY_MAX_ROWS]
        rows = [[o["name"], o["type"], o["slot_count"],
                 trim(o["description"], DESC_CHARS)] for o in shown]
        deck.table(slide, left, top, width, height,
                   ["Object", "Type", "Slots", "Description"], rows,
                   widths=[0.26, 0.22, 0.08, 0.44], size=11)
        if len(objects) > len(shown):
            deck.text_block(slide, left, top + height, width, Inches(0.3),
                            [f"{len(objects) - len(shown)} more objects not "
                             f"listed."], size=10, bullet=False)
        return slide, [f"Source: {self.model_file}, object inventory. "
                       f"{len(shown)} of {len(objects)} objects listed."]

    # -- policy ------------------------------------------------------------
    def _policy(self, spec):
        deck = self.deck
        refs = spec.get("refs") or {}
        chosen = (self.sets[refs["set"]] if refs.get("set")
                  else next(iter(self.sets.values())))
        slide = deck.add_slide("title_only", self._title_for(
            spec, f"Operating policy: {chosen['name']}"))
        left, top, width, height = deck.content_box()
        top, height = self._bullets_under(slide, spec, left, top, width, height)

        agenda = chosen["agenda"]
        caption = [f"Rule set \"{chosen['name']}\", agenda order {agenda}"
                   + (" -- the bottom rule fires first, and a rule listed "
                      "higher fires later and overrides it."
                      if agenda == "ASCENDING" else ".")]
        if chosen["description"]:
            caption.append(trim(chosen["description"], 180))
        cap_h = Inches(0.3) * len(caption)
        deck.text_block(slide, left, top, width, cap_h, caption, size=12,
                        bullet=False)
        top += cap_h + Inches(0.08)
        height -= cap_h + Inches(0.08)

        # Agenda positions are counted over the whole set, so a slide showing
        # three groups still carries the numbers the modeler sees in RiverWare.
        groups = [g for g in chosen["groups"] if g["kind"] == "POLICY_GROUP"]
        priority = {}
        position = 0
        for group in groups:
            for item in group["items"]:
                position += 1
                priority[id(item)] = position
        wanted = refs.get("groups")
        if wanted:
            # The spec's order wins: for an ASCENDING set the file order is the
            # reverse of the firing order, and which one a room should see is
            # the spec author's call.
            by_name = {g["name"]: g for g in groups}
            groups = [by_name[n] for n in wanted]

        annotations = spec.get("annotations") or {}
        rows, total = [], 0
        for group in groups:
            flag = "" if group["active"] else "  [inactive]"
            total += 1
            if len(rows) < POLICY_MAX_ROWS:
                rows.append(["", f"{group['name']}{flag}",
                             trim(annotations.get(group["name"], ""),
                                  NOTE_CHARS)])
            for item in group["items"]:
                total += 1
                if len(rows) >= POLICY_MAX_ROWS:
                    continue
                note = annotations.get(f"{group['name']}/{item['name']}") \
                    or item["description"]
                mark = "" if item["active"] else "  [inactive]"
                rows.append([priority[id(item)], f"    {item['name']}{mark}",
                             trim(note, NOTE_CHARS)])
        deck.table(slide, left, top, width, height,
                   ["#", "Group and rules", "What it does"], rows,
                   widths=[0.06, 0.36, 0.58], size=11)
        if total > len(rows):
            deck.text_block(slide, left, top + height, width, Inches(0.3),
                            [f"{total - len(rows)} more rows not listed."],
                            size=10, bullet=False)
        return slide, [f"Source: rule set \"{chosen['name']}\" as stored in "
                       f"{self.model_file}. Numbers are agenda positions, not "
                       f"execution order for {agenda} sets."]

    # -- series ------------------------------------------------------------
    def _series(self, spec):
        deck = self.deck
        refs = spec.get("refs") or {}
        picks = [self.series[r] for r in refs["series"]]
        slide = deck.add_slide("title_only", self._title_for(
            spec, " and ".join(sorted({p["slot"] for p in picks}))))
        left, top, width, height = deck.content_box()
        top, height = self._bullets_under(slide, spec, left, top, width, height)

        axis, columns, stride = align_series(picks)
        chart_data = CategoryChartData()
        chart_data.categories = axis
        chart_data.categories.number_format = "yyyy-mm-dd"
        units = {p["unit"] for p in picks}
        for pick, values in zip(picks, columns):
            label = f"{pick['object']} {pick['slot']}"
            if len(units) > 1:
                label += f" ({pick['unit']})"
            chart_data.add_series(label, values, number_format="General")
        frame = slide.shapes.add_chart(XL_CHART_TYPE.LINE, left, top, width,
                                       height, chart_data)
        chart = frame.chart
        chart.has_title = False
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.font.size = Pt(11)
        for i, plot_series in enumerate(chart.plots[0].series):
            plot_series.smooth = False
            plot_series.format.line.color.rgb = RGBColor.from_string(
                SERIES_COLORS[i % len(SERIES_COLORS)])
            plot_series.format.line.width = Pt(1.75)
        if len(units) == 1:
            axis_obj = chart.value_axis
            axis_obj.has_title = True
            axis_obj.axis_title.text_frame.text = next(iter(units))
            axis_obj.axis_title.text_frame.paragraphs[0].font.size = Pt(11)
        else:
            deck.warnings.append(
                "series slide mixes units (" + ", ".join(sorted(units))
                + "); one value axis cannot label them all")

        span = (f"{axis[0]:%Y-%m-%d} to {axis[-1]:%Y-%m-%d}" if axis else "empty")
        thinned = (f" Plotted every {stride} timesteps to keep the chart light."
                   if stride > 1 else "")
        return slide, [
            f"Source: {self.model_file}, stored results for "
            + ", ".join(f"{p['object']}.{p['slot']} ({p['unit']})" for p in picks)
            + f". {span}, {len(axis)} points.{thinned} Gaps are timesteps with "
              f"no stored value."]

    # -- table -------------------------------------------------------------
    def _table(self, spec):
        deck = self.deck
        refs = spec.get("refs") or {}
        data = self.tables[refs["table"]]
        slide = deck.add_slide("title_only", self._title_for(
            spec, f"{data['object']} {data['slot']}"))
        left, top, width, height = deck.content_box()
        top, height = self._bullets_under(slide, spec, left, top, width, height)

        cols = data["cols"] or [f"Column {i + 1}"
                                for i in range(len(data["rows"][0]))]
        cols = [c or f"Column {i + 1}" for i, c in enumerate(cols)]
        shown = data["rows"][:TABLE_MAX_ROWS]
        deck.table(slide, left, top, width, height, cols,
                   [[fmt_num(v) for v in row] for row in shown], size=12)
        if len(data["rows"]) > len(shown):
            deck.text_block(slide, left, top + height, width, Inches(0.3),
                            [f"{len(data['rows']) - len(shown)} more rows not "
                             f"listed."], size=10, bullet=False)
        return slide, [f"Source: {self.model_file}, {refs['table']} "
                       f"({data['kind']}), {len(data['rows'])} rows. "
                       f"{len(shown)} shown."]

    # -- chart-xy ----------------------------------------------------------
    def _chart_xy(self, spec):
        deck = self.deck
        refs = spec.get("refs") or {}
        data = self.tables[refs["table"]]
        cols = data["cols"]
        x_name = refs.get("x") or cols[0]
        x_index = cols.index(x_name)
        y_names = refs.get("y") or [c for i, c in enumerate(cols)
                                    if i != x_index]
        slide = deck.add_slide("title_only", self._title_for(
            spec, f"{data['object']} {data['slot']}"))
        left, top, width, height = deck.content_box()
        top, height = self._bullets_under(slide, spec, left, top, width, height)

        chart_data = XyChartData()
        for y_name in y_names:
            y_index = cols.index(y_name)
            series = chart_data.add_series(y_name)
            points = [(row[x_index], row[y_index]) for row in data["rows"]
                      if row[x_index] is not None and row[y_index] is not None]
            for x_value, y_value in sorted(points):
                series.add_data_point(x_value, y_value)
        frame = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS, left, top, width, height,
            chart_data)
        chart = frame.chart
        chart.has_title = False
        chart.has_legend = len(y_names) > 1
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        chart.font.size = Pt(11)
        for i, plot_series in enumerate(chart.plots[0].series):
            plot_series.smooth = False
            plot_series.format.line.color.rgb = RGBColor.from_string(
                SERIES_COLORS[i % len(SERIES_COLORS)])
            plot_series.format.line.width = Pt(1.75)
        chart.category_axis.has_title = True
        chart.category_axis.axis_title.text_frame.text = x_name
        chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(11)
        if len(y_names) == 1:
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.text = y_names[0]
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(11)
        return slide, [f"Source: {self.model_file}, {refs['table']}, "
                       f"{len(data['rows'])} rows plotted as {x_name} against "
                       + ", ".join(y_names) + "."]

    # -- bullets and caveats ----------------------------------------------
    def _bullets(self, spec):
        deck = self.deck
        slide = deck.add_slide("title_only", self._title_for(spec, ""))
        left, top, width, height = deck.content_box()
        lines = [str(b) for b in spec["bullets"]]
        deck.text_block(slide, left, top, width, height, lines,
                        size=18 if len(lines) <= BULLET_MAX else 15)
        return slide, []

    def _caveats(self, spec):
        deck = self.deck
        slide = deck.add_slide("title_only",
                               self._title_for(spec, "Reading this deck"))
        left, top, width, height = deck.content_box()
        lines = CAVEAT_LINES + [str(b) for b in (spec.get("bullets") or [])]
        deck.text_block(slide, left, top, width, height, lines, size=15)
        return slide, [f"Source: {self.model_file}."]


# --------------------------------------------------------------------- output
def build_deck(digest: dict, spec: dict, when: str,
               template: Path | None) -> tuple[Deck, list[str]]:
    footer = f"{digest['model']['file']}  |  {when}"
    deck = Deck(template, footer)
    renderer = Renderer(deck, digest, spec, when)
    for slide_spec in spec["slides"]:
        renderer.render(slide_spec)
    props = deck.prs.core_properties
    props.title = spec.get("deck_title") or digest["model"]["file"]
    props.author = ""
    props.last_modified_by = ""
    props.revision = 1
    props.created = FIXED_TIMESTAMP
    props.modified = FIXED_TIMESTAMP
    return deck, deck.warnings


def _rezip(blob: bytes, rewrite=None) -> bytes:
    """Rebuild a zip with fixed entry dates, optionally editing each member."""
    out = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(blob)) as src, \
            zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as dst:
        for info in src.infolist():
            data = src.read(info.filename)
            if rewrite is not None:
                data = rewrite(info.filename, data)
            entry = zipfile.ZipInfo(info.filename, date_time=FIXED_ZIP_DATE)
            entry.compress_type = zipfile.ZIP_DEFLATED
            entry.external_attr = info.external_attr
            entry.create_system = 0
            dst.writestr(entry, data)
    return out.getvalue()


def _fixed_member(name: str, data: bytes) -> bytes:
    """Strip the clock out of one package member.

    A chart carries its data as an embedded Excel workbook, which is a zip of
    its own with its own timestamps and its own creation date. Left alone it
    makes every render of the same chart a different file.
    """
    if name.endswith(".xlsx"):
        return _rezip(data, _fixed_member)
    if name.endswith("docProps/core.xml"):
        return re.sub(rb">\d{4}-\d{2}-\d{2}T[\d:.]+Z?<",
                      FIXED_TIMESTAMP.strftime(">%Y-%m-%dT%H:%M:%SZ<").encode(),
                      data)
    return data


def save_deterministic(prs, out_path: Path) -> None:
    """Write the package with the clock taken out of it.

    python-pptx stamps every zip entry with the current time, so two renders
    of one spec would differ in bytes while being the same document.
    """
    buffer = io.BytesIO()
    prs.save(buffer)
    out_path.write_bytes(_rezip(buffer.getvalue(), _fixed_member))


# ----------------------------------------------------------------------- main
def load_spec(path: Path) -> dict:
    data = json.loads(path.read_text(encoding="utf-8"))
    if isinstance(data, list):
        data = {"slides": data}
    if not isinstance(data, dict):
        raise ValueError("spec must be a JSON object with a 'slides' list")
    return data


def parse_args(argv: list[str]) -> dict:
    opts: dict = {"spec": None, "auto": False, "template": None, "out": None,
                  "rls": None, "date": None, "model": None}
    positional: list[str] = []
    i = 0
    while i < len(argv):
        arg = argv[i]
        if arg == "--auto":
            opts["auto"] = True
        elif arg in ("--spec", "--template", "--rls", "--date", "-o", "--output"):
            if i + 1 >= len(argv):
                raise ValueError(f"{arg} needs a value")
            key = {"-o": "out", "--output": "out"}.get(arg, arg.lstrip("-"))
            opts[key] = argv[i + 1]
            i += 1
        elif arg.startswith("-"):
            raise ValueError(f"unknown option {arg}")
        else:
            positional.append(arg)
        i += 1
    if len(positional) != 1:
        raise ValueError("expected exactly one model file")
    opts["model"] = positional[0]
    return opts


def main(argv: list[str]) -> int:
    try:
        opts = parse_args(argv)
    except ValueError as exc:
        if argv:
            print(f"error: {exc}", file=sys.stderr)
        print(__doc__)
        return 1
    if bool(opts["spec"]) == opts["auto"]:
        print("error: pass exactly one of --spec <deck.json> or --auto",
              file=sys.stderr)
        return 1

    model = Path(opts["model"])
    if not model.exists() or model.suffix.lower() != ".mdl":
        print(f"error: {model} is not an existing .mdl file", file=sys.stderr)
        return 2
    if PPTX_ERROR is not None:
        print("error: this script needs the python-pptx package.\n"
              "  pip install python-pptx\n"
              f"  (import failed: {PPTX_ERROR})", file=sys.stderr)
        return 2

    digest = build_digest(model, include_policy=True)
    if opts["rls"]:
        rls = Path(opts["rls"])
        if not rls.exists() or rls.suffix.lower() != ".rls":
            print(f"error: {rls} is not an existing .rls file", file=sys.stderr)
            return 2
        supplied = policy_from_rls(rls.read_text(encoding="utf-8",
                                                 errors="replace"))
        # A supplied ruleset is the operating policy the modeler wants shown,
        # so it leads; sets serialized in the model stay available by name.
        digest["policy"]["sets"] = (supplied["sets"]
                                    + digest["policy"].get("sets", []))

    when = opts["date"] or _date.today().isoformat()
    if not re.fullmatch(r"\d{4}-\d{2}-\d{2}", when):
        print(f"error: --date {when} is not YYYY-MM-DD", file=sys.stderr)
        return 2

    out = Path(opts["out"]) if opts["out"] \
        else model.with_name(model.stem + "_deck.pptx")
    spec_out = out.with_suffix(".json")

    if opts["auto"]:
        spec = auto_spec(digest, model.stem, when)
    else:
        spec_path = Path(opts["spec"])
        if not spec_path.exists():
            print(f"error: {spec_path} not found", file=sys.stderr)
            return 2
        try:
            spec = load_spec(spec_path)
        except (ValueError, json.JSONDecodeError) as exc:
            print(f"error: {spec_path}: {exc}", file=sys.stderr)
            return 2
        if spec.get("date"):
            when = spec["date"]
        if spec.get("template") and not opts["template"]:
            opts["template"] = spec["template"]

    errs = validate_spec(spec, digest)
    if errs:
        print("error: the deck spec was rejected; nothing was written:",
              file=sys.stderr)
        for err in errs:
            print(f"  - {err}", file=sys.stderr)
        return 3

    template = Path(opts["template"]) if opts["template"] else None
    if template is not None and not template.exists():
        print(f"error: template {template} not found", file=sys.stderr)
        return 2

    deck, warnings = build_deck(digest, spec, when, template)
    save_deterministic(deck.prs, out)
    if opts["auto"]:
        spec_out.write_text(json.dumps(spec, indent=2) + "\n", encoding="utf-8")

    kinds: dict[str, int] = {}
    for slide_spec in spec["slides"]:
        kinds[slide_spec["type"]] = kinds.get(slide_spec["type"], 0) + 1
    print(f"build_pptx.py: {model.name} -> {out.name}")
    print(f"  slides:  {len(spec['slides'])} ("
          + ", ".join(f"{n} {k}" for k, n in sorted(kinds.items())) + ")")
    if opts["auto"]:
        print(f"  spec:    {spec_out.name} (edit it and re-run with --spec)")
    if template is not None:
        print(f"  template: {template.name}")
    files = digest.get("policy", {}).get("referenced_files", [])
    if files and not digest["policy"]["sets"]:
        print("  ruleset: none stored in the model. It records these paths, "
              "which were not opened:")
        for name in files:
            print(f"    - {name}")
    for warning in warnings:
        print(f"  warning: {warning}")
    print("\nOnly PowerPoint validates a .pptx. Open the deck and check the "
          "schematic, the charts and the speaker notes before presenting it.")
    return 0


if __name__ == "__main__":
    try:
        _code = main(sys.argv[1:])
        sys.stdout.flush()  # surface a closed pipe here, not at shutdown
    except BrokenPipeError:
        os.dup2(os.open(os.devnull, os.O_WRONLY), sys.stdout.fileno())
        _code = 0
    raise SystemExit(_code)
