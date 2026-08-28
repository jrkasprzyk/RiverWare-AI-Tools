"""Renderer tests for skills/present-riverware-model/build_pptx.py (TEST-001..007).

Run from the repo root:  python -m unittest discover -s tests

Most tests run against tests/fixtures/mini_model.mdl, the same miniature the
annotate tests use: three objects and a two-group rule set, with no stored
results and no populated lookup table. That absence is itself under test -- a
model nobody has run is the common case, and the deck has to say so instead of
drawing an empty chart.

The slide types that need real data (series, table, chart-xy) run against
examples/ArborBasin/ArborBasin.mdl, whose digest is parsed once and shared.

Everything here is skipped when python-pptx is missing, which is the same
condition under which the script refuses to run.
"""
import json
import subprocess
import sys
import tempfile
import unittest
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path

REPO = Path(__file__).resolve().parent.parent
BUILD = REPO / "skills" / "present-riverware-model" / "build_pptx.py"
FIXTURE = REPO / "tests" / "fixtures" / "mini_model.mdl"
ARBOR = REPO / "examples" / "ArborBasin" / "ArborBasin.mdl"
DASHBOARD = REPO / "skills" / "visualize-riverware-model" / "digest_to_json.py"

sys.path.insert(0, str(REPO / "skills" / "present-riverware-model"))
sys.path.insert(0, str(REPO / "skills" / "visualize-riverware-model"))
import build_pptx  # noqa: E402
from digest_to_json import build_digest, layout_nodes  # noqa: E402

HAVE_PPTX = build_pptx.PPTX_ERROR is None
needs_pptx = unittest.skipUnless(
    HAVE_PPTX, f"python-pptx not installed ({build_pptx.PPTX_ERROR})")

FIXED_DATE = "2026-08-28"


def run_build(*args) -> subprocess.CompletedProcess:
    return subprocess.run([sys.executable, str(BUILD), *map(str, args)],
                          capture_output=True, text=True, cwd=str(REPO))


class DeckCase(unittest.TestCase):
    """Base: render a spec into a temp deck and hand back the path."""

    _digests: dict = {}

    @classmethod
    def digest(cls, model: Path) -> dict:
        if model not in cls._digests:
            cls._digests[model] = build_digest(model, include_policy=True)
        return cls._digests[model]

    def render(self, spec: dict, model: Path = FIXTURE, *extra):
        """Returns (completed_process, deck_path_or_None) inside a temp dir."""
        tmp = Path(tempfile.mkdtemp())
        self.addCleanup(lambda: None)  # temp dirs are left to the OS on Windows
        spec_path = tmp / "deck.json"
        spec_path.write_text(json.dumps(spec), encoding="utf-8")
        out = tmp / "deck.pptx"
        proc = run_build(model, "--spec", spec_path, "-o", out,
                         "--date", FIXED_DATE, *extra)
        return proc, (out if out.exists() else None)

    def render_ok(self, spec: dict, model: Path = FIXTURE, *extra) -> Path:
        proc, out = self.render(spec, model, *extra)
        self.assertEqual(proc.returncode, 0,
                         f"stdout:\n{proc.stdout}\nstderr:\n{proc.stderr}")
        self.assertIsNotNone(out)
        return out


# ------------------------------------------------------------------ TEST-001
class TestSpecValidation(DeckCase):
    """A wrong spec is reported in full and produces no file at all."""

    def test_unknown_slide_type_lists_the_vocabulary(self):
        errs = build_pptx.validate_spec(
            {"slides": [{"type": "histogram"}]}, self.digest(FIXTURE))
        self.assertEqual(len(errs), 1)
        self.assertIn("histogram", errs[0])
        for kind in build_pptx.SLIDE_TYPES:
            self.assertIn(kind, errs[0])

    def test_unresolvable_refs_name_the_valid_targets(self):
        digest = self.digest(ARBOR)
        errs = build_pptx.validate_spec({"slides": [
            {"type": "series", "refs": {"series": ["Nowhere.Pool Elevation"]}},
            {"type": "table", "refs": {"table": "Aspen.Nope"}},
            {"type": "network", "refs": {"objects": ["Ghost"]}},
        ]}, digest)
        self.assertEqual(len(errs), 3)
        self.assertIn("Aspen.Pool Elevation", errs[0])
        self.assertIn("Aspen.Elevation Volume Table", errs[1])
        self.assertIn("Aspen", errs[2])

    def test_unknown_policy_group_is_rejected(self):
        errs = build_pptx.validate_spec({"slides": [
            {"type": "policy", "refs": {"set": "Mini Set",
                                        "groups": ["Nowhere Rules"]}}]},
            self.digest(FIXTURE))
        self.assertEqual(len(errs), 1)
        self.assertIn("Cora Rules", errs[0])

    def test_duplicate_ids_are_rejected(self):
        errs = build_pptx.validate_spec({"slides": [
            {"type": "title", "id": "a"}, {"type": "caveats", "id": "a"}]},
            self.digest(FIXTURE))
        self.assertEqual(len(errs), 1)
        self.assertIn("duplicate id", errs[0])

    def test_empty_spec_is_rejected(self):
        self.assertTrue(build_pptx.validate_spec({"slides": []},
                                                 self.digest(FIXTURE)))

    @needs_pptx
    def test_a_rejected_spec_writes_nothing(self):
        proc, out = self.render({"slides": [{"type": "histogram", "id": "h"},
                                            {"type": "title", "id": "t"}]})
        self.assertEqual(proc.returncode, 3)
        self.assertIsNone(out, "a rejected spec must not leave a partial deck")
        self.assertIn("nothing was written", proc.stderr)


# ------------------------------------------------------------------ TEST-002
@needs_pptx
class TestDeterminism(DeckCase):
    """REQ-006: one spec plus one model is one deck, on every run."""

    def test_two_runs_are_byte_identical(self):
        spec = {"deck_title": "Determinism", "slides": [
            {"type": "title", "id": "t"},
            {"type": "network", "id": "n"},
            {"type": "policy", "id": "p"},
            {"type": "caveats", "id": "c"}]}
        first = self.render_ok(spec).read_bytes()
        second = self.render_ok(spec).read_bytes()
        self.assertEqual(first, second)

    def test_charts_are_deterministic_too(self):
        spec = {"slides": [{"type": "series", "id": "s", "refs": {
            "series": ["Aspen.Pool Elevation", "Birch.Pool Elevation"]}}]}
        self.assertEqual(self.render_ok(spec, ARBOR).read_bytes(),
                         self.render_ok(spec, ARBOR).read_bytes())


# ------------------------------------------------------------------ TEST-003
@needs_pptx
class TestSlideVocabulary(DeckCase):
    """Every slide type renders, and what it emits is a valid package."""

    ARBOR_SPEC = {"deck_title": "Every slide type", "slides": [
        {"type": "title", "id": "cover", "bullets": ["A story, not a dump."]},
        {"type": "bullets", "id": "agenda", "title": "Agenda",
         "bullets": ["System", "Policy", "Results"]},
        {"type": "network", "id": "net",
         "refs": {"objects": ["Aspen", "Aspen to Birch", "Birch"]}},
        {"type": "summary", "id": "sum", "refs": {"types": ["StorageReservoir"]}},
        {"type": "policy", "id": "pol",
         "refs": {"set": "Arbor Basin Rules (from MRM run)",
                  "groups": ["Aspen Rules"]},
         "annotations": {"Aspen Rules": "Holds Aspen on its guide curve."}},
        {"type": "series", "id": "ser",
         "refs": {"series": ["Aspen.Pool Elevation"]}},
        {"type": "table", "id": "tab",
         "refs": {"table": "Aspen.Elevation Volume Table"}},
        {"type": "chart-xy", "id": "xy",
         "refs": {"table": "Aspen.Elevation Volume Table", "x": "Storage",
                  "y": ["Pool Elevation"]}},
        {"type": "caveats", "id": "end"}]}

    @classmethod
    def setUpClass(cls):
        if not HAVE_PPTX:
            raise unittest.SkipTest("python-pptx not installed")
        cls.deck = DeckCase.render_ok(cls("run"), cls.ARBOR_SPEC, ARBOR)

    def test_slide_count_matches_the_spec(self):
        from pptx import Presentation
        prs = Presentation(str(self.deck))
        self.assertEqual(len(prs.slides), len(self.ARBOR_SPEC["slides"]))

    def test_every_emitted_part_is_well_formed_xml(self):
        with zipfile.ZipFile(self.deck) as zf:
            names = zf.namelist()
            for name in names:
                if name.endswith((".xml", ".rels")):
                    ET.fromstring(zf.read(name))   # raises on malformed XML
        self.assertIn("[Content_Types].xml", names)
        self.assertIn("ppt/presentation.xml", names)
        charts = [n for n in names if n.startswith("ppt/charts/chart")]
        self.assertEqual(len(charts), 2, "one series chart and one xy chart")

    def test_the_zip_inventory_matches_the_content_types(self):
        with zipfile.ZipFile(self.deck) as zf:
            declared = ET.fromstring(zf.read("[Content_Types].xml"))
            overrides = {el.get("PartName").lstrip("/") for el in declared
                         if el.tag.endswith("Override")}
            names = set(zf.namelist())
        self.assertTrue(overrides <= names,
                        f"declared but absent: {sorted(overrides - names)}")

    def test_time_series_uses_a_date_axis(self):
        """RISK-001: a date axis is what keeps a 366-point chart legible."""
        with zipfile.ZipFile(self.deck) as zf:
            charts = [zf.read(n).decode("utf-8") for n in zf.namelist()
                      if n.startswith("ppt/charts/chart")]
        self.assertTrue(any("c:dateAx" in c for c in charts))

    def test_nothing_is_drawn_off_the_slide(self):
        from pptx import Presentation
        prs = Presentation(str(self.deck))
        for index, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.left is None or shape.width is None:
                    continue
                with self.subTest(slide=index, shape=shape.name):
                    self.assertGreaterEqual(shape.left, 0)
                    self.assertGreaterEqual(shape.top, 0)
                    self.assertLessEqual(shape.left + shape.width,
                                         prs.slide_width)
                    self.assertLessEqual(shape.top + shape.height,
                                         prs.slide_height)

    def test_every_data_slide_carries_provenance_notes(self):
        """REQ-011: notes name the source, so a forwarded deck stays checkable."""
        from pptx import Presentation
        prs = Presentation(str(self.deck))
        for index, (slide, spec) in enumerate(
                zip(prs.slides, self.ARBOR_SPEC["slides"]), 1):
            if spec["type"] == "bullets":
                continue
            with self.subTest(slide=index, type=spec["type"]):
                notes = slide.notes_slide.notes_text_frame.text
                self.assertIn("ArborBasin.mdl", notes)

    def test_the_footer_names_the_model_and_the_date_only(self):
        """REQ-012: a staleness signal, not an attribution line."""
        from pptx import Presentation
        prs = Presentation(str(self.deck))
        texts = [shape.text_frame.text for shape in prs.slides[0].shapes
                 if shape.has_text_frame]
        footer = [t for t in texts if FIXED_DATE in t]
        self.assertTrue(footer)
        self.assertIn("ArborBasin.mdl", footer[0])
        for banned in ("AI", "Claude", "generated by", "python-pptx"):
            self.assertNotIn(banned, footer[0])


# ------------------------------------------------------------------ TEST-004
@needs_pptx
class TestPaging(DeckCase):
    """Long tables are cut to a readable slide and say what was left out."""

    def test_table_slide_reports_the_rows_it_dropped(self):
        from pptx import Presentation
        deck = self.render_ok({"slides": [{"type": "table", "id": "t", "refs": {
            "table": "Aspen.Elevation Volume Table"}}]}, ARBOR)
        prs = Presentation(str(deck))
        slide = prs.slides[0]
        rows = [sh.table for sh in slide.shapes if sh.has_table][0]
        self.assertEqual(len(rows.rows), build_pptx.TABLE_MAX_ROWS + 1)
        text = " ".join(sh.text_frame.text for sh in slide.shapes
                        if sh.has_text_frame)
        digest = self.digest(ARBOR)
        total = [t for t in digest["tables"]
                 if t["object"] == "Aspen"
                 and t["slot"] == "Elevation Volume Table"][0]
        self.assertIn(f"{len(total['rows']) - build_pptx.TABLE_MAX_ROWS} more "
                      f"rows", text)

    def test_summary_slide_reports_the_objects_it_dropped(self):
        from pptx import Presentation
        deck = self.render_ok({"slides": [{"type": "summary", "id": "s"}]},
                              ARBOR)
        text = " ".join(sh.text_frame.text for sh in
                        Presentation(str(deck)).slides[0].shapes
                        if sh.has_text_frame)
        self.assertIn("more objects not listed", text)


# ------------------------------------------------------------------ TEST-005
@needs_pptx
class TestModelWithoutResults(DeckCase):
    """A model nobody has run must produce an honest deck, not an empty chart."""

    def test_auto_skips_series_slides(self):
        tmp = Path(tempfile.mkdtemp())
        out = tmp / "mini_deck.pptx"
        proc = run_build(FIXTURE, "--auto", "-o", out, "--date", FIXED_DATE)
        self.assertEqual(proc.returncode, 0, proc.stderr)
        spec = json.loads((tmp / "mini_deck.json").read_text(encoding="utf-8"))
        kinds = [s["type"] for s in spec["slides"]]
        self.assertNotIn("series", kinds)
        self.assertEqual(kinds[:3], ["title", "network", "summary"])
        self.assertEqual(kinds[-1], "caveats")

    def test_a_requested_series_is_an_error_not_a_blank_chart(self):
        proc, out = self.render({"slides": [{"type": "series", "id": "s",
                                             "refs": {"series": ["Cora.Outflow"]}}]})
        self.assertEqual(proc.returncode, 3)
        self.assertIsNone(out)
        self.assertIn("stores no results", proc.stderr)

    def test_the_summary_says_when_no_ruleset_is_available(self):
        """REQ-010: thinner content is stated, never silent."""
        from pptx import Presentation
        digest = self.digest(FIXTURE)
        spec = {"slides": [{"type": "summary", "id": "s"}]}
        stripped = dict(digest, policy={"sets": [], "referenced_files": []})
        deck = build_pptx.Deck(None, "footer")
        renderer = build_pptx.Renderer(deck, stripped, spec, FIXED_DATE)
        renderer.render(spec["slides"][0])
        text = " ".join(sh.text_frame.text for sh in deck.prs.slides[0].shapes
                        if sh.has_text_frame)
        self.assertIn("Ruleset not included", text)

    def test_a_policy_slide_without_a_ruleset_names_the_rls_option(self):
        errs = build_pptx.validate_spec(
            {"slides": [{"type": "policy", "id": "p"}]},
            {"objects": [], "series": [], "tables": [],
             "policy": {"sets": []}})
        self.assertEqual(len(errs), 1)
        self.assertIn("--rls", errs[0])


# ------------------------------------------------------------------ TEST-006
@needs_pptx
class TestTemplate(DeckCase):
    """REQ-007: a client template is honoured where it can be, never fatally."""

    SPEC = {"slides": [{"type": "title", "id": "t"},
                       {"type": "network", "id": "n"},
                       {"type": "caveats", "id": "c"}]}

    ACCENT = b"BB2244"

    def _template(self, keep_layouts) -> Path:
        """A template with a recognisable accent colour and chosen layouts.

        The accent is stamped into the theme part directly: what matters is
        that the colour survives into the rendered deck, and reaching for it
        through the object model would test python-pptx, not this renderer.
        """
        import io
        import re
        from pptx import Presentation
        prs = Presentation()
        if keep_layouts is not None:
            layouts = prs.slide_masters[0].slide_layouts
            for layout in list(layouts):
                if layout.name not in keep_layouts:
                    layouts.remove(layout)
        buffer = io.BytesIO()
        prs.save(buffer)
        path = Path(tempfile.mkdtemp()) / "client.pptx"
        with zipfile.ZipFile(io.BytesIO(buffer.getvalue())) as src, \
                zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as dst:
            for info in src.infolist():
                data = src.read(info.filename)
                if info.filename == "ppt/theme/theme1.xml":
                    data = re.sub(rb'(<a:accent1><a:srgbClr val=")[0-9A-F]{6}',
                                  rb"\g<1>" + self.ACCENT, data)
                dst.writestr(info, data)
        return path

    def test_theme_is_inherited_when_layouts_are_present(self):
        from pptx import Presentation
        template = self._template(None)
        with zipfile.ZipFile(template) as zf:
            self.assertIn(self.ACCENT, zf.read("ppt/theme/theme1.xml"))
        deck = self.render_ok(self.SPEC, FIXTURE, "--template", template)
        with zipfile.ZipFile(deck) as zf:
            self.assertIn(self.ACCENT, zf.read("ppt/theme/theme1.xml"))
        self.assertEqual(len(Presentation(str(deck)).slides), 3)

    def test_a_template_without_the_named_layouts_still_renders(self):
        proc, out = self.render(self.SPEC, FIXTURE, "--template",
                                self._template({"Blank"}))
        self.assertEqual(proc.returncode, 0, proc.stderr)
        self.assertIsNotNone(out)
        self.assertIn("warning:", proc.stdout)
        self.assertIn("layout", proc.stdout)

    def test_a_missing_template_is_a_clean_error(self):
        proc, out = self.render(self.SPEC, FIXTURE, "--template",
                                "no_such_template.pptx")
        self.assertEqual(proc.returncode, 2)
        self.assertIsNone(out)
        self.assertIn("template", proc.stderr)


# ------------------------------------------------------------------ TEST-007
class TestSharedDigest(unittest.TestCase):
    """The dashboard must be unaffected by the extraction the deck skill added."""

    def test_the_dashboard_payload_carries_no_policy_tree(self):
        proc = subprocess.run([sys.executable, str(DASHBOARD), str(FIXTURE)],
                              capture_output=True, text=True, cwd=str(REPO))
        self.assertEqual(proc.returncode, 0, proc.stderr)
        digest = json.loads(proc.stdout)
        self.assertEqual(sorted(digest),
                         ["links", "model", "objects", "series", "tables"])

    def test_html_output_never_embeds_the_policy_tree(self):
        with tempfile.TemporaryDirectory() as tmp:
            out = Path(tmp) / "dash.html"
            proc = subprocess.run(
                [sys.executable, str(DASHBOARD), str(FIXTURE), "--html",
                 "--policy", "-o", str(out)],
                capture_output=True, text=True, cwd=str(REPO))
            self.assertEqual(proc.returncode, 0, proc.stderr)
            self.assertNotIn('"policy"', out.read_text(encoding="utf-8"))

    def test_build_digest_adds_policy_only_when_asked(self):
        plain = build_digest(FIXTURE)
        self.assertNotIn("policy", plain)
        with_policy = build_digest(FIXTURE, include_policy=True)
        names = [s["name"] for s in with_policy["policy"]["sets"]]
        self.assertEqual(names, ["Mini Set"])
        groups = with_policy["policy"]["sets"][0]["groups"]
        self.assertEqual([g["name"] for g in groups],
                         ["Cora Rules", "Roberto Rules"])

    def test_layout_places_every_object_once(self):
        digest = build_digest(FIXTURE)
        layout = layout_nodes(digest["objects"], digest["links"])
        self.assertEqual(sorted(n["name"] for n in layout["nodes"]),
                         sorted(o["name"] for o in digest["objects"]))
        for node in layout["nodes"]:
            self.assertLessEqual(node["x"] + layout["node_width"],
                                 layout["width"])
            self.assertLessEqual(node["y"] + layout["node_height"],
                                 layout["height"])

    def test_recorded_ruleset_paths_are_reported_not_opened(self):
        digest = build_digest(ARBOR, include_policy=True)
        self.assertIn("/ModelsAndData/MRM/ArborBasin_Rules.rls",
                      digest["policy"]["referenced_files"])


# --------------------------------------------------------------------- CON-002
@needs_pptx
class TestConsoleOutput(DeckCase):
    def test_stdout_is_ascii(self):
        tmp = Path(tempfile.mkdtemp())
        proc = run_build(ARBOR, "--auto", "-o", tmp / "d.pptx",
                         "--date", FIXED_DATE)
        self.assertEqual(proc.returncode, 0, proc.stderr)
        proc.stdout.encode("ascii")

    def test_usage_is_printed_without_arguments(self):
        proc = run_build()
        self.assertEqual(proc.returncode, 1)
        self.assertIn("--spec", proc.stdout)

    def test_spec_and_auto_are_mutually_exclusive(self):
        proc = run_build(FIXTURE, "--auto", "--spec", "x.json")
        self.assertEqual(proc.returncode, 1)
        self.assertIn("exactly one", proc.stderr)


if __name__ == "__main__":
    unittest.main()
